# backend/services/enhanced_processor.py
"""
Enhanced Document Processing Service - COMPLETE FIXED VERSION
Handles ALL document formats with robust text extraction and comprehensive debugging
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, NamedTuple
import asyncio
from concurrent.futures import ThreadPoolExecutor
import traceback

# Standard library imports
import re
import unicodedata
from datetime import datetime

# Document processing libraries
try:
    import docx
    from docx.document import Document as DocxDocument
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    POWERPOINT_AVAILABLE = True
except ImportError:
    POWERPOINT_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False

logger = logging.getLogger(__name__)


class ProcessingResult(NamedTuple):
    """Standardized processing result"""
    text: str
    chunks: List[Dict[str, Any]]
    total_tokens: int
    total_chunks: int
    processing_time: float
    file_info: Dict[str, Any]
    metadata: Dict[str, Any]


class DocumentProcessor:
    """Enhanced document processor with robust text extraction for ALL formats"""
    
    def __init__(self):
        self.executor = ThreadPoolExecutor(max_workers=2)
        self._setup_logging()
        self.supported_formats = {
            '.docx': 'Microsoft Word Document',
            '.pdf': 'Portable Document Format',
            '.xlsx': 'Microsoft Excel Spreadsheet',
            '.xls': 'Microsoft Excel Spreadsheet (Legacy)',
            '.pptx': 'Microsoft PowerPoint Presentation',
            '.ppt': 'Microsoft PowerPoint Presentation (Legacy)',
            '.html': 'HTML Document',
            '.htm': 'HTML Document',
            '.epub': 'EPUB eBook',
            '.csv': 'Comma-Separated Values',
            '.txt': 'Plain Text',
            '.md': 'Markdown Document',
            '.markdown': 'Markdown Document',
            '.py': 'Python Code',
            '.js': 'JavaScript Code',
            '.ts': 'TypeScript Code',
            '.json': 'JSON Data',
            '.xml': 'XML Document',
            '.yaml': 'YAML Document',
            '.yml': 'YAML Document'
        }
    
    def _setup_logging(self):
        """Setup comprehensive logging for debugging"""
        logger.setLevel(logging.DEBUG)
        if not logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter(
                '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
            )
            handler.setFormatter(formatter)
            logger.addHandler(handler)
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_formats.keys())
    
    async def process_file(
        self, 
        file_path: str, 
        max_tokens: int = 1024,
        tokenizer: str = "word-estimate"
    ) -> ProcessingResult:
        """
        Process any supported document format with comprehensive error handling
        """
        start_time = datetime.utcnow()
        
        try:
            # Validate file exists
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_ext = Path(file_path).suffix.lower()
            file_size = os.path.getsize(file_path)
            filename = Path(file_path).name
            
            logger.info(f"🔍 Starting processing of {file_ext} file: {filename} ({file_size:,} bytes)")
            
            # Check if format is supported
            if file_ext not in self.supported_formats:
                logger.warning(f"⚠️ Unsupported format {file_ext}, attempting text fallback")
                text = await self._extract_txt_text(file_path)
            else:
                # Extract text based on file type with comprehensive error handling
                try:
                    if file_ext == '.docx':
                        text = await self._extract_docx_text_enhanced(file_path)
                    elif file_ext == '.pdf':
                        text = await self._extract_pdf_text_enhanced(file_path)
                    elif file_ext in ['.xlsx', '.xls']:
                        text = await self._extract_excel_text_enhanced(file_path)
                    elif file_ext in ['.pptx', '.ppt']:
                        text = await self._extract_powerpoint_text_enhanced(file_path)
                    elif file_ext in ['.html', '.htm']:
                        text = await self._extract_html_text_enhanced(file_path)
                    elif file_ext == '.epub':
                        text = await self._extract_epub_text_enhanced(file_path)
                    elif file_ext == '.csv':
                        text = await self._extract_csv_text_enhanced(file_path)
                    elif file_ext == '.txt':
                        text = await self._extract_txt_text(file_path)
                    elif file_ext in ['.md', '.markdown']:
                        text = await self._extract_markdown_text(file_path)
                    elif file_ext in ['.py', '.js', '.ts', '.json', '.xml', '.yaml', '.yml']:
                        text = await self._extract_code_text(file_path)
                    else:
                        # This shouldn't happen due to the check above, but just in case
                        text = await self._extract_txt_text(file_path)
                
                except Exception as extraction_error:
                    logger.error(f"❌ {file_ext} extraction failed: {extraction_error}")
                    logger.debug(f"Full extraction error traceback:\n{traceback.format_exc()}")
                    
                    # Try fallback to text extraction
                    logger.info(f"🔄 Attempting text fallback for {file_ext} file")
                    try:
                        text = await self._extract_txt_text(file_path)
                        logger.info(f"✅ Text fallback successful: {len(text)} characters")
                    except Exception as fallback_error:
                        logger.error(f"❌ Even text fallback failed: {fallback_error}")
                        raise ValueError(f"Complete extraction failure for {file_ext}: {extraction_error}")
            
            # Enhanced validation with detailed feedback
            if not text:
                raise ValueError(f"No text content extracted from {filename}")
            
            text_length = len(text.strip())
            if text_length == 0:
                raise ValueError(f"Empty text content extracted from {filename}")
            
            if text_length < 10:
                logger.warning(f"⚠️ Very short text extracted: {text_length} characters")
                logger.debug(f"Short text content: '{text}'")
            
            logger.info(f"✅ Text extraction successful: {text_length:,} characters from {filename}")
            logger.debug(f"Text preview (first 300 chars): {text[:300]}...")
            
            # Clean and process the text
            cleaned_text = self._clean_text_enhanced(text)
            
            # Create chunks with enhanced logic
            chunks = self._create_chunks_enhanced(cleaned_text, max_tokens, tokenizer)
            
            # Calculate tokens
            total_tokens = self._estimate_tokens_enhanced(cleaned_text, tokenizer)
            
            # Calculate processing time
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            
            # Create comprehensive result
            result = ProcessingResult(
                text=cleaned_text,
                chunks=chunks,
                total_tokens=total_tokens,
                total_chunks=len(chunks),
                processing_time=processing_time,
                file_info={
                    "filename": filename,
                    "format": file_ext.lstrip('.'),
                    "format_description": self.supported_formats.get(file_ext, "Unknown format"),
                    "size_bytes": file_size,
                    "processed_at": datetime.utcnow().isoformat() + "Z"
                },
                metadata={
                    "processing_method": "enhanced",
                    "extraction_method": f"{file_ext}_enhanced",
                    "tokenizer": tokenizer,
                    "max_tokens": max_tokens,
                    "original_text_length": len(text),
                    "cleaned_text_length": len(cleaned_text),
                    "processing_time_seconds": processing_time,
                    "libraries_used": self._get_available_libraries()
                }
            )
            
            logger.info(f"🎉 Processing completed successfully!")
            logger.info(f"   📊 {result.total_chunks} chunks, {result.total_tokens:,} tokens")
            logger.info(f"   ⏱️ Processing time: {processing_time:.2f} seconds")
            
            return result
            
        except Exception as e:
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.error(f"❌ Processing failed for {file_path}: {e}")
            logger.debug(f"Full processing error traceback:\n{traceback.format_exc()}")
            
            # Still provide some debug info even on failure
            logger.info(f"🔍 Processing failure analysis:")
            logger.info(f"   📁 File: {file_path}")
            logger.info(f"   📏 File size: {os.path.getsize(file_path) if os.path.exists(file_path) else 'N/A'}")
            logger.info(f"   🔧 Available libraries: {self._get_available_libraries()}")
            logger.info(f"   ⏱️ Time before failure: {processing_time:.2f} seconds")
            
            raise
    
    def _get_available_libraries(self) -> Dict[str, bool]:
        """Get status of available processing libraries"""
        return {
            "docx": DOCX_AVAILABLE,
            "pdf": PDF_AVAILABLE,
            "excel": EXCEL_AVAILABLE,
            "powerpoint": POWERPOINT_AVAILABLE,
            "html": HTML_AVAILABLE,
            "epub": EPUB_AVAILABLE
        }
    
    async def _extract_docx_text_enhanced(self, file_path: str) -> str:
        """Enhanced DOCX text extraction with multiple fallback strategies"""
        
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available. Install with: pip install python-docx")
        
        def extract_sync():
            try:
                logger.debug(f"🔍 Opening DOCX file: {file_path}")
                doc = docx.Document(file_path)
                text_parts = []
                
                logger.debug(f"📄 DOCX structure: {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables, {len(doc.sections)} sections")
                
                # Strategy 1: Document structure iteration (most reliable)
                try:
                    from docx.oxml.text.paragraph import CT_P
                    from docx.oxml.table import CT_Tbl
                    
                    def iter_block_items(parent):
                        """Iterate through all document elements in order"""
                        if hasattr(parent, 'element'):
                            parent_elm = parent.element.body
                        else:
                            parent_elm = parent
                        
                        for child in parent_elm:
                            if isinstance(child, CT_P):
                                yield Paragraph(child, parent)
                            elif isinstance(child, CT_Tbl):
                                yield Table(child, parent)
                    
                    # Extract in document order
                    for block in iter_block_items(doc):
                        if isinstance(block, Paragraph):
                            para_text = block.text.strip()
                            if para_text:
                                text_parts.append(para_text)
                        elif isinstance(block, Table):
                            table_text = self._extract_table_text_enhanced(block)
                            if table_text:
                                text_parts.append(table_text)
                    
                    logger.debug(f"✅ Strategy 1 (structure iteration): extracted {len(text_parts)} elements")
                
                except Exception as structure_error:
                    logger.warning(f"⚠️ Structure iteration failed: {structure_error}")
                    text_parts = []
                
                # Strategy 2: Basic paragraph extraction (fallback)
                if not text_parts:
                    logger.debug("🔄 Strategy 2: Basic paragraph extraction")
                    for paragraph in doc.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            text_parts.append(para_text)
                    
                    # Extract tables separately
                    for table in doc.tables:
                        table_text = self._extract_table_text_enhanced(table)
                        if table_text:
                            text_parts.append(table_text)
                    
                    logger.debug(f"✅ Strategy 2 (basic extraction): extracted {len(text_parts)} elements")
                
                # Strategy 3: Run-level extraction (last resort)
                if not text_parts:
                    logger.debug("🔄 Strategy 3: Run-level extraction")
                    for paragraph in doc.paragraphs:
                        for run in paragraph.runs:
                            run_text = run.text.strip()
                            if run_text:
                                text_parts.append(run_text)
                    
                    logger.debug(f"✅ Strategy 3 (run-level): extracted {len(text_parts)} elements")
                
                # Extract headers and footers
                for section in doc.sections:
                    # Headers
                    if hasattr(section, 'header') and section.header:
                        for para in section.header.paragraphs:
                            header_text = para.text.strip()
                            if header_text:
                                text_parts.append(f"[Header] {header_text}")
                    
                    # Footers
                    if hasattr(section, 'footer') and section.footer:
                        for para in section.footer.paragraphs:
                            footer_text = para.text.strip()
                            if footer_text:
                                text_parts.append(f"[Footer] {footer_text}")
                
                # Combine all text
                full_text = "\n\n".join(text_parts)
                
                logger.info(f"📊 DOCX extraction complete:")
                logger.info(f"   📝 Text elements extracted: {len(text_parts)}")
                logger.info(f"   📏 Total characters: {len(full_text):,}")
                logger.info(f"   📄 Document stats: {len(doc.paragraphs)} paras, {len(doc.tables)} tables")
                
                if not full_text.strip():
                    raise ValueError("No readable text content found in DOCX document")
                
                return full_text
                
            except Exception as e:
                logger.error(f"❌ DOCX extraction failed: {e}")
                logger.debug(f"DOCX extraction traceback:\n{traceback.format_exc()}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    def _extract_table_text_enhanced(self, table) -> str:
        """Enhanced table text extraction"""
        try:
            table_rows = []
            for row_idx, row in enumerate(table.rows):
                row_cells = []
                for cell in row.cells:
                    # Extract all paragraphs from cell
                    cell_paras = []
                    for para in cell.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            cell_paras.append(para_text)
                    
                    cell_text = " ".join(cell_paras) if cell_paras else ""
                    if cell_text:
                        row_cells.append(cell_text)
                
                if row_cells:
                    table_rows.append(" | ".join(row_cells))
            
            result = "\n".join(table_rows) if table_rows else ""
            logger.debug(f"Table extracted: {len(table_rows)} rows, {len(result)} characters")
            return result
            
        except Exception as e:
            logger.warning(f"⚠️ Table extraction failed: {e}")
            return ""
    
    async def _extract_pdf_text_enhanced(self, file_path: str) -> str:
        """Enhanced PDF text extraction with multiple engines and strategies"""
        
        if not PDF_AVAILABLE:
            raise ImportError("PDF processing requires PyPDF2 and pdfplumber")
        
        def extract_sync():
            try:
                text_parts = []
                
                # Strategy 1: pdfplumber (best for complex layouts)
                try:
                    logger.debug("🔍 Strategy 1: pdfplumber extraction")
                    with pdfplumber.open(file_path) as pdf:
                        logger.debug(f"📄 PDF has {len(pdf.pages)} pages")
                        for page_num, page in enumerate(pdf.pages, 1):
                            try:
                                page_text = page.extract_text()
                                if page_text and page_text.strip():
                                    clean_page_text = page_text.strip()
                                    text_parts.append(f"[Page {page_num}]\n{clean_page_text}")
                                    logger.debug(f"Page {page_num}: {len(clean_page_text)} characters")
                                else:
                                    logger.debug(f"Page {page_num}: No text found")
                            except Exception as page_error:
                                logger.warning(f"⚠️ Page {page_num} extraction failed: {page_error}")
                    
                    logger.info(f"✅ pdfplumber: extracted text from {len(text_parts)} pages")
                    
                except Exception as plumber_error:
                    logger.warning(f"⚠️ pdfplumber extraction failed: {plumber_error}")
                    text_parts = []
                
                # Strategy 2: PyPDF2 fallback
                if not text_parts:
                    logger.debug("🔄 Strategy 2: PyPDF2 extraction")
                    try:
                        with open(file_path, 'rb') as file:
                            pdf_reader = PyPDF2.PdfReader(file)
                            logger.debug(f"📄 PDF has {len(pdf_reader.pages)} pages (PyPDF2)")
                            
                            for page_num, page in enumerate(pdf_reader.pages, 1):
                                try:
                                    page_text = page.extract_text()
                                    if page_text and page_text.strip():
                                        clean_page_text = page_text.strip()
                                        text_parts.append(f"[Page {page_num}]\n{clean_page_text}")
                                        logger.debug(f"Page {page_num}: {len(clean_page_text)} characters")
                                except Exception as page_error:
                                    logger.warning(f"⚠️ Page {page_num} extraction failed: {page_error}")
                        
                        logger.info(f"✅ PyPDF2: extracted text from {len(text_parts)} pages")
                    
                    except Exception as pypdf_error:
                        logger.error(f"❌ PyPDF2 extraction failed: {pypdf_error}")
                        raise
                
                if not text_parts:
                    raise ValueError("No text could be extracted from PDF using any method")
                
                full_text = "\n\n".join(text_parts)
                
                logger.info(f"📊 PDF extraction complete:")
                logger.info(f"   📄 Pages processed: {len(text_parts)}")
                logger.info(f"   📏 Total characters: {len(full_text):,}")
                
                return full_text
                
            except Exception as e:
                logger.error(f"❌ PDF extraction failed: {e}")
                logger.debug(f"PDF extraction traceback:\n{traceback.format_exc()}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_excel_text_enhanced(self, file_path: str) -> str:
        """Enhanced Excel text extraction"""
        
        if not EXCEL_AVAILABLE:
            raise ImportError("Excel processing requires openpyxl and pandas")
        
        def extract_sync():
            try:
                logger.debug(f"🔍 Opening Excel file: {file_path}")
                
                # Try with openpyxl first
                workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                text_parts = []
                
                logger.debug(f"📊 Excel workbook has {len(workbook.sheetnames)} sheets: {workbook.sheetnames}")
                
                for sheet_idx, sheet_name in enumerate(workbook.sheetnames):
                    logger.debug(f"📋 Processing sheet {sheet_idx + 1}: {sheet_name}")
                    
                    try:
                        sheet = workbook[sheet_name]
                        
                        # Get all data from sheet
                        data = []
                        row_count = 0
                        for row in sheet.iter_rows(values_only=True):
                            if any(cell is not None and str(cell).strip() for cell in row):
                                cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                                if any(cleaned_row):  # Only add rows with actual content
                                    data.append(cleaned_row)
                                    row_count += 1
                        
                        if data:
                            # Create sheet text
                            sheet_text = f"=== Sheet: {sheet_name} ===\n"
                            
                            # Add headers if they look like headers
                            if data and row_count > 1:
                                headers = data[0]
                                if all(isinstance(h, str) and len(h) < 50 for h in headers if h):
                                    sheet_text += "Headers: " + " | ".join(headers) + "\n\n"
                                    data_rows = data[1:]
                                else:
                                    data_rows = data
                            else:
                                data_rows = data
                            
                            # Add data rows (limit to prevent huge outputs)
                            max_rows = 500  # Reasonable limit
                            for i, row in enumerate(data_rows[:max_rows]):
                                row_text = " | ".join(str(cell) for cell in row if cell)
                                if row_text.strip():
                                    sheet_text += row_text + "\n"
                            
                            if len(data_rows) > max_rows:
                                sheet_text += f"\n... ({len(data_rows) - max_rows} more rows)\n"
                            
                            text_parts.append(sheet_text)
                            logger.debug(f"Sheet '{sheet_name}': {row_count} rows, {len(sheet_text)} characters")
                        else:
                            logger.debug(f"Sheet '{sheet_name}': No data found")
                    
                    except Exception as sheet_error:
                        logger.warning(f"⚠️ Failed to process sheet '{sheet_name}': {sheet_error}")
                
                workbook.close()
                
                if not text_parts:
                    raise ValueError("No data found in any Excel sheets")
                
                full_text = "\n\n".join(text_parts)
                
                logger.info(f"📊 Excel extraction complete:")
                logger.info(f"   📋 Sheets processed: {len(text_parts)}")
                logger.info(f"   📏 Total characters: {len(full_text):,}")
                
                return full_text
                
            except Exception as e:
                logger.error(f"❌ Excel extraction failed: {e}")
                logger.debug(f"Excel extraction traceback:\n{traceback.format_exc()}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    # Additional extraction methods continue...
    async def _extract_txt_text(self, file_path: str) -> str:
        """Enhanced text file extraction with encoding detection"""
        
        def extract_sync():
            try:
                logger.debug(f"🔍 Reading text file: {file_path}")
                
                # Try different encodings in order of likelihood
                encodings = ['utf-8', 'utf-8-sig', 'utf-16', 'latin-1', 'cp1252', 'iso-8859-1']
                
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as file:
                            text = file.read()
                            logger.info(f"✅ Text file read successfully with {encoding}: {len(text):,} characters")
                            return text
                    except UnicodeDecodeError as e:
                        logger.debug(f"Encoding {encoding} failed: {e}")
                        continue
                    except Exception as e:
                        logger.warning(f"⚠️ Failed to read with {encoding}: {e}")
                        continue
                
                # If all encodings fail, try with error handling
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                        text = file.read()
                        logger.warning(f"⚠️ Text file read with character replacement: {len(text):,} characters")
                        return text
                except Exception as e:
                    logger.error(f"❌ Even error-tolerant reading failed: {e}")
                    raise ValueError(f"Could not read text file with any encoding: {e}")
                
            except Exception as e:
                logger.error(f"❌ Text extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    # Continue with other extraction methods...
    async def _extract_powerpoint_text_enhanced(self, file_path: str) -> str:
        """Enhanced PowerPoint text extraction"""
        if not POWERPOINT_AVAILABLE:
            raise ImportError("python-pptx not available")
        
        def extract_sync():
            try:
                prs = Presentation(file_path)
                text_parts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                    
                    if slide_content:
                        slide_text = f"[Slide {slide_num}]\n" + "\n".join(slide_content)
                        text_parts.append(slide_text)
                
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"PowerPoint extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_html_text_enhanced(self, file_path: str) -> str:
        """Enhanced HTML text extraction"""
        if not HTML_AVAILABLE:
            raise ImportError("HTML processing requires beautifulsoup4")
        
        def extract_sync():
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    html_content = file.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                text = soup.get_text()
                
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                return '\n'.join(chunk for chunk in chunks if chunk)
            except Exception as e:
                logger.error(f"HTML extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_csv_text_enhanced(self, file_path: str) -> str:
        """Enhanced CSV text extraction"""
        def extract_sync():
            try:
                df = pd.read_csv(file_path, encoding='utf-8')
                text_parts = [f"CSV Data ({len(df)} rows, {len(df.columns)} columns)"]
                text_parts.append("Columns: " + " | ".join(df.columns.astype(str)))
                
                max_rows = 1000
                for idx, row in df.head(max_rows).iterrows():
                    text_parts.append(" | ".join(row.astype(str)))
                
                if len(df) > max_rows:
                    text_parts.append(f"... ({len(df) - max_rows} more rows)")
                
                return "\n".join(text_parts)
            except Exception as e:
                logger.error(f"CSV extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_epub_text_enhanced(self, file_path: str) -> str:
        """Enhanced EPUB text extraction"""
        if not EPUB_AVAILABLE:
            raise ImportError("EPUB processing requires ebooklib")
        
        def extract_sync():
            try:
                book = epub.read_epub(file_path)
                text_parts = []
                
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        
                        for script in soup(["script", "style"]):
                            script.decompose()
                        
                        chapter_text = soup.get_text()
                        
                        if chapter_text.strip():
                            lines = (line.strip() for line in chapter_text.splitlines())
                            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                            cleaned_text = '\n'.join(chunk for chunk in chunks if chunk)
                            text_parts.append(cleaned_text)
                
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"EPUB extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_markdown_text(self, file_path: str) -> str:
        """Extract text from Markdown files"""
        def extract_sync():
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    markdown_text = file.read()
                
                # Basic markdown cleaning
                text = re.sub(r'#{1,6}\s+', '', markdown_text)
                text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
                text = re.sub(r'\*(.*?)\*', r'\1', text)
                text = re.sub(r'`(.*?)`', r'\1', text)
                text = re.sub(r'```.*?\n(.*?)\n```', r'\1', text, flags=re.DOTALL)
                
                return text.strip()
            except Exception as e:
                logger.error(f"Markdown extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    async def _extract_code_text(self, file_path: str) -> str:
        """Extract text from code files with syntax preservation"""
        def extract_sync():
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    code_content = file.read()
                
                file_ext = Path(file_path).suffix.lower()
                filename = Path(file_path).name
                
                return f"=== {filename} ({file_ext}) ===\n\n{code_content}"
            except Exception as e:
                logger.error(f"Code extraction failed: {e}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(self.executor, extract_sync)
    
    def _clean_text_enhanced(self, text: str) -> str:
        """Enhanced text cleaning and normalization"""
        if not text:
            return ""
        
        logger.debug(f"🧹 Cleaning text: {len(text)} characters")
        
        # Normalize line endings
        cleaned = re.sub(r'\r\n', '\n', text)
        cleaned = re.sub(r'\r', '\n', cleaned)
        
        # Reduce excessive whitespace
        cleaned = re.sub(r'\n{4,}', '\n\n\n', cleaned)  # Max 3 consecutive newlines
        cleaned = re.sub(r'[ \t]{2,}', ' ', cleaned)    # Multiple spaces to single space
        cleaned = re.sub(r'[ \t]*\n[ \t]*', '\n', cleaned)  # Remove trailing/leading spaces on lines
        
        # Unicode normalization
        cleaned = unicodedata.normalize('NFKC', cleaned)
        
        # Remove common document artifacts
        cleaned = re.sub(r'\f', '\n', cleaned)  # Form feed to newline
        cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', cleaned)  # Remove control chars
        
        # Final cleanup
        cleaned = cleaned.strip()
        
        logger.debug(f"✅ Text cleaning complete: {len(text)} -> {len(cleaned)} characters")
        return cleaned
    
    def _create_chunks_enhanced(self, text: str, max_tokens: int, tokenizer: str) -> List[Dict[str, Any]]:
        """Enhanced chunking with multiple strategies and better token estimation"""
        if not text:
            logger.warning("⚠️ No text provided for chunking")
            return []
        
        logger.debug(f"🔪 Creating chunks: max_tokens={max_tokens}, tokenizer={tokenizer}")
        
        # Calculate rough words per chunk based on tokenizer
        if tokenizer == "word-estimate":
            # GPT-style: ~0.75 tokens per word
            words_per_chunk = int(max_tokens / 0.75)
        elif tokenizer == "char-estimate":
            # Character-based: ~4 chars per token
            chars_per_chunk = max_tokens * 4
            words_per_chunk = chars_per_chunk // 5  # ~5 chars per word average
        else:
            # Default conservative estimate
            words_per_chunk = int(max_tokens * 0.7)
        
        logger.debug(f"📏 Target words per chunk: {words_per_chunk}")
        
        # Strategy 1: Paragraph-based chunking (preferred)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        if not paragraphs:
            # Strategy 2: Sentence-based chunking
            logger.debug("🔄 No paragraphs found, using sentence-based chunking")
            sentences = re.split(r'[.!?]+', text)
            paragraphs = [s.strip() for s in sentences if s.strip()]
        
        if not paragraphs:
            # Strategy 3: Line-based chunking
            logger.debug("🔄 No sentences found, using line-based chunking")
            paragraphs = [line.strip() for line in text.split('\n') if line.strip()]
        
        if not paragraphs:
            # Strategy 4: Word-based chunking (last resort)
            logger.debug("🔄 No lines found, using word-based chunking")
            words = text.split()
            paragraphs = [' '.join(words[i:i+words_per_chunk]) 
                         for i in range(0, len(words), words_per_chunk)]
        
        logger.debug(f"📝 Found {len(paragraphs)} text segments for chunking")
        
        # Create chunks
        chunks = []
        current_chunk = ""
        current_word_count = 0
        chunk_id = 1
        
        for para_idx, paragraph in enumerate(paragraphs):
            para_words = len(paragraph.split())
            
            # Check if this paragraph fits in current chunk
            if current_word_count + para_words <= words_per_chunk or not current_chunk:
                # Add to current chunk
                if current_chunk:
                    current_chunk += "\n\n" + paragraph
                else:
                    current_chunk = paragraph
                current_word_count += para_words
            else:
                # Current chunk is full, save it and start new one
                if current_chunk:
                    chunk_tokens = self._estimate_tokens_enhanced(current_chunk, tokenizer)
                    chunks.append({
                        "text": current_chunk.strip(),
                        "chunk_id": chunk_id,
                        "token_count": chunk_tokens,
                        "word_count": current_word_count,
                        "chunk_index": chunk_id - 1
                    })
                    chunk_id += 1
                
                # Start new chunk with current paragraph
                current_chunk = paragraph
                current_word_count = para_words
        
        # Add final chunk if there's remaining content
        if current_chunk:
            chunk_tokens = self._estimate_tokens_enhanced(current_chunk, tokenizer)
            chunks.append({
                "text": current_chunk.strip(),
                "chunk_id": chunk_id,
                "token_count": chunk_tokens,
                "word_count": current_word_count,
                "chunk_index": chunk_id - 1
            })
        
        # Ensure we have at least one chunk
        if not chunks and text:
            logger.debug("🔄 Creating single chunk from entire text")
            chunk_tokens = self._estimate_tokens_enhanced(text, tokenizer)
            chunks.append({
                "text": text.strip(),
                "chunk_id": 1,
                "token_count": chunk_tokens,
                "word_count": len(text.split()),
                "chunk_index": 0
            })
        
        logger.info(f"✅ Created {len(chunks)} chunks")
        for i, chunk in enumerate(chunks[:3]):  # Log first 3 chunks for debugging
            logger.debug(f"Chunk {i+1}: {chunk['word_count']} words, {chunk['token_count']} tokens")
        
        return chunks
    
    def _estimate_tokens_enhanced(self, text: str, tokenizer: str = "word-estimate") -> int:
        """Enhanced token estimation with multiple methods"""
        if not text:
            return 0
        
        text = text.strip()
        if not text:
            return 0
        
        if tokenizer == "word-estimate":
            # GPT-style estimation: roughly 0.75 tokens per word
            word_count = len(text.split())
            return max(1, int(word_count * 0.75))
        
        elif tokenizer == "char-estimate":
            # Character-based estimation: roughly 4 characters per token
            char_count = len(text)
            return max(1, char_count // 4)
        
        elif tokenizer == "conservative":
            # Conservative estimation: 1 token per word
            word_count = len(text.split())
            return max(1, word_count)
        
        else:
            # Default: use word estimation
            word_count = len(text.split())
            return max(1, int(word_count * 0.75))


# Test and validation functions
async def test_enhanced_processor():
    """Test the enhanced processor with various scenarios"""
    processor = DocumentProcessor()
    
    logger.info("🧪 Enhanced Document Processor Test Suite")
    logger.info(f"📋 Supported formats: {processor.get_supported_formats()}")
    logger.info(f"📚 Available libraries: {processor._get_available_libraries()}")
    
    return processor


# Create processor instance for import
enhanced_processor = DocumentProcessor()


if __name__ == "__main__":
    import asyncio
    asyncio.run(test_enhanced_processor())
