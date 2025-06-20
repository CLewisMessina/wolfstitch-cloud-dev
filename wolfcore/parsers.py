"""
Wolfcore File Parsers - Multi-format Document Processing
Extracted from Wolfstitch Desktop App: processing/extract.py

Supports 40+ file formats with intelligent text extraction:
- Documents: PDF, DOCX, TXT, EPUB, HTML, MD
- Code: Python, JavaScript, Java, C++, Go, Rust, etc.
- Data: CSV, JSON, YAML, XML
- Presentations: PPTX
- Spreadsheets: XLSX

Key Features:
- Automatic encoding detection
- Format-specific optimization  
- Metadata preservation
- Error resilience with fallbacks
- Progress reporting for large files
"""

import os
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List, Union, BinaryIO
from dataclasses import dataclass, field
from datetime import datetime
import chardet
import tempfile

from .exceptions import ParsingError, UnsupportedFormatError

logger = logging.getLogger(__name__)


@dataclass
class ParsedFile:
    """Container for parsed file content and metadata"""
    
    # Core content
    text: str
    filename: str
    format: str
    
    # Metadata
    size_bytes: int = 0
    encoding: Optional[str] = None
    language: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    line_count: Optional[int] = None
    
    # Processing info
    extraction_method: str = "unknown"
    processing_time: float = 0.0
    warnings: List[str] = field(default_factory=list)
    
    # Additional metadata (format-specific)
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def __post_init__(self):
        """Calculate derived metrics after initialization"""
        if self.text:
            self.word_count = len(self.text.split())
            self.line_count = self.text.count('\n') + 1


class FileParser:
    """
    Multi-format file parser with progressive enhancement
    Extracted and enhanced from desktop app's extract.py
    """
    
    # Supported formats mapping
    SUPPORTED_FORMATS = {
        # Documents
        '.pdf': 'pdf',
        '.docx': 'docx', 
        '.doc': 'doc',
        '.txt': 'txt',
        '.epub': 'epub',
        '.html': 'html',
        '.htm': 'html',
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.rtf': 'rtf',
        
        # Spreadsheets and data
        '.xlsx': 'xlsx',
        '.csv': 'csv',
        '.json': 'json',
        '.yaml': 'yaml',
        '.yml': 'yaml',
        '.xml': 'xml',
        '.toml': 'toml',
        '.ini': 'ini',
        
        # Presentations
        '.pptx': 'pptx',
        
        # Code files
        '.py': 'python',
        '.js': 'javascript',
        '.jsx': 'javascript',
        '.ts': 'typescript', 
        '.tsx': 'typescript',
        '.java': 'java',
        '.cpp': 'cpp',
        '.cc': 'cpp',
        '.cxx': 'cpp',
        '.c': 'c',
        '.h': 'c_header',
        '.hpp': 'cpp_header',
        '.go': 'go',
        '.rs': 'rust',
        '.rb': 'ruby',
        '.php': 'php',
        '.swift': 'swift',
        '.kt': 'kotlin',
        '.cs': 'csharp',
        '.r': 'r',
        '.scala': 'scala',
        '.sh': 'shell',
        '.bash': 'shell',
        '.zsh': 'shell',
        '.fish': 'shell',
    }
    
    def __init__(self, max_file_size_mb: int = 100):
        """
        Initialize file parser
        
        Args:
            max_file_size_mb: Maximum file size to process (MB)
        """
        self.max_file_size_bytes = max_file_size_mb * 1024 * 1024
        self.temp_files = []  # Track temp files for cleanup
        
    def parse(self, file_path: Union[str, Path, BinaryIO], 
              filename: Optional[str] = None) -> ParsedFile:
        """
        Parse file and extract text content
        
        Args:
            file_path: Path to file or file-like object
            filename: Override filename (useful for file-like objects)
            
        Returns:
            ParsedFile with extracted content and metadata
            
        Raises:
            ParsingError: If parsing fails
            UnsupportedFormatError: If format not supported
        """
        start_time = datetime.now()
        
        try:
            # Handle different input types
            if hasattr(file_path, 'read'):
                # File-like object
                return self._parse_file_object(file_path, filename)
            else:
                # File path
                file_path = Path(file_path)
                return self._parse_file_path(file_path)
                
        except Exception as e:
            logger.error(f"Parsing failed for {filename or file_path}: {e}")
            raise ParsingError(f"Failed to parse file: {e}") from e
        finally:
            # Cleanup any temporary files
            self._cleanup_temp_files()
            
    def _parse_file_path(self, file_path: Path) -> ParsedFile:
        """Parse file from filesystem path"""
        
        # Validate file exists and size
        if not file_path.exists():
            raise ParsingError(f"File not found: {file_path}")
            
        file_size = file_path.stat().st_size
        if file_size > self.max_file_size_bytes:
            raise ParsingError(
                f"File too large: {file_size / 1024 / 1024:.1f}MB "
                f"(max: {self.max_file_size_bytes / 1024 / 1024}MB)"
            )
        
        # Determine format
        file_format = self._detect_format(file_path.name)
        
        # Extract content based on format
        start_time = datetime.now()
        
        if file_format == 'pdf':
            result = self._parse_pdf(file_path)
        elif file_format == 'docx':
            result = self._parse_docx(file_path)
        elif file_format == 'epub':
            result = self._parse_epub(file_path)
        elif file_format == 'xlsx':
            result = self._parse_xlsx(file_path)
        elif file_format == 'pptx':
            result = self._parse_pptx(file_path)
        elif file_format in ['txt', 'python', 'javascript', 'typescript', 
                            'java', 'cpp', 'c', 'go', 'rust', 'ruby', 
                            'php', 'swift', 'kotlin', 'csharp', 'r', 
                            'scala', 'shell', 'markdown']:
            result = self._parse_text_file(file_path, file_format)
        elif file_format in ['html', 'xml']:
            result = self._parse_markup(file_path, file_format)
        elif file_format in ['json', 'yaml', 'toml', 'ini', 'csv']:
            result = self._parse_structured_data(file_path, file_format)
        else:
            raise UnsupportedFormatError(f"Unsupported format: {file_format}")
        
        # Set metadata
        result.filename = file_path.name
        result.format = file_format
        result.size_bytes = file_size
        result.processing_time = (datetime.now() - start_time).total_seconds()
        
        return result
        
    def _parse_file_object(self, file_obj: BinaryIO, filename: str) -> ParsedFile:
        """Parse file from file-like object"""
        
        # Save to temporary file for processing
        with tempfile.NamedTemporaryFile(delete=False, 
                                       suffix=f".{filename.split('.')[-1]}") as tmp:
            content = file_obj.read()
            tmp.write(content)
            tmp_path = Path(tmp.name)
            self.temp_files.append(tmp_path)
        
        # Parse the temporary file
        result = self._parse_file_path(tmp_path)
        result.filename = filename
        
        return result
        
    def _detect_format(self, filename: str) -> str:
        """Detect file format from filename"""
        suffix = Path(filename).suffix.lower()
        
        if suffix in self.SUPPORTED_FORMATS:
            return self.SUPPORTED_FORMATS[suffix]
        else:
            raise UnsupportedFormatError(f"Unsupported file extension: {suffix}")
            
#stitch point
    def _parse_pdf(self, file_path: Path) -> ParsedFile:
            """Extract text from PDF using pdfminer.six"""
            try:
                from pdfminer.high_level import extract_text
                from pdfminer.pdfpage import PDFPage
            
                # Extract text content
                text = extract_text(str(file_path))
                
                # Count pages
                with open(file_path, 'rb') as file:
                    page_count = len(list(PDFPage.get_pages(file)))
                
                return ParsedFile(
                    text=text,
                    filename="",  # Will be set by caller
                    format="pdf",
                    extraction_method="pdfminer.six",
                    metadata={"page_count": page_count}
                )
                
            except ImportError:
                raise ParsingError("PDF parsing requires pdfminer.six: pip install pdfminer.six")
            except Exception as e:
                raise ParsingError(f"PDF parsing failed: {e}")
    
    def _parse_docx(self, file_path: Path) -> ParsedFile:
        """Extract text from DOCX using python-docx"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract text from tables
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
            
            # Combine all text
            all_text = paragraphs + table_text
            text = "\n\n".join(all_text)
            
            # Extract metadata
            props = doc.core_properties
            metadata = {
                "author": props.author or "Unknown",
                "created": str(props.created) if props.created else None,
                "modified": str(props.modified) if props.modified else None,
                "title": props.title or "",
                "subject": props.subject or "",
                "paragraph_count": len(paragraphs),
                "table_count": len(doc.tables)
            }
            
            return ParsedFile(
                text=text,
                filename="",
                format="docx",
                extraction_method="python-docx",
                metadata=metadata
            )
            
        except ImportError:
            raise ParsingError("DOCX parsing requires python-docx: pip install python-docx")
        except Exception as e:
            raise ParsingError(f"DOCX parsing failed: {e}")
    
    def _parse_epub(self, file_path: Path) -> ParsedFile:
        """Extract text from EPUB using BeautifulSoup"""
        try:
            import zipfile
            from bs4 import BeautifulSoup
            
            chapters = []
            
            with zipfile.ZipFile(file_path, 'r') as epub:
                # Find content files
                for file_info in epub.filelist:
                    if file_info.filename.endswith(('.html', '.xhtml', '.htm')):
                        try:
                            content = epub.read(file_info.filename).decode('utf-8')
                            soup = BeautifulSoup(content, 'html.parser')
                            
                            # Remove script and style elements
                            for script in soup(["script", "style"]):
                                script.decompose()
                            
                            # Extract text
                            text = soup.get_text()
                            
                            # Clean up whitespace
                            lines = (line.strip() for line in text.splitlines())
                            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                            text = ' '.join(chunk for chunk in chunks if chunk)
                            
                            if text.strip():
                                chapters.append(text)
                                
                        except Exception as e:
                            logger.warning(f"Failed to extract from {file_info.filename}: {e}")
            
            full_text = "\n\n".join(chapters)
            
            return ParsedFile(
                text=full_text,
                filename="",
                format="epub",
                extraction_method="zipfile + BeautifulSoup",
                metadata={"chapter_count": len(chapters)}
            )
            
        except ImportError:
            raise ParsingError("EPUB parsing requires BeautifulSoup: pip install beautifulsoup4")
        except Exception as e:
            raise ParsingError(f"EPUB parsing failed: {e}")
    
    def _parse_xlsx(self, file_path: Path) -> ParsedFile:
        """Extract text from Excel files using openpyxl"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, data_only=True)
            
            all_text = []
            sheet_info = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                # Extract cell values
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                if sheet_text:
                    all_text.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_text))
                    sheet_info[sheet_name] = len(sheet_text)
            
            full_text = "\n\n".join(all_text)
            
            return ParsedFile(
                text=full_text,
                filename="",
                format="xlsx",
                extraction_method="openpyxl",
                metadata={
                    "sheet_count": len(workbook.sheetnames),
                    "sheets": sheet_info
                }
            )
            
        except ImportError:
            raise ParsingError("XLSX parsing requires openpyxl: pip install openpyxl")
        except Exception as e:
            raise ParsingError(f"XLSX parsing failed: {e}")
    
    def _parse_pptx(self, file_path: Path) -> ParsedFile:
        """Extract text from PowerPoint files using python-pptx"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slide_texts = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_content = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slide_text = f"Slide {i}:\n" + "\n".join(slide_content)
                    slide_texts.append(slide_text)
            
            full_text = "\n\n".join(slide_texts)
            
            return ParsedFile(
                text=full_text,
                filename="",
                format="pptx",
                extraction_method="python-pptx",
                metadata={"slide_count": len(prs.slides)}
            )
            
        except ImportError:
            raise ParsingError("PPTX parsing requires python-pptx: pip install python-pptx")
        except Exception as e:
            raise ParsingError(f"PPTX parsing failed: {e}")
    
    def _parse_text_file(self, file_path: Path, file_format: str) -> ParsedFile:
        """Parse plain text and code files with encoding detection"""
        
        # Detect encoding
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            
        encoding_info = chardet.detect(raw_data)
        encoding = encoding_info.get('encoding', 'utf-8')
        confidence = encoding_info.get('confidence', 0.0)
        
        warnings = []
        if confidence < 0.8:
            warnings.append(f"Low encoding confidence: {confidence:.2f}")
        
        try:
            # Try detected encoding
            text = raw_data.decode(encoding)
        except (UnicodeDecodeError, TypeError):
            # Fallback to utf-8 with error handling
            try:
                text = raw_data.decode('utf-8', errors='replace')
                encoding = 'utf-8'
                warnings.append("Used UTF-8 fallback encoding")
            except:
                # Last resort - latin-1 (never fails)
                text = raw_data.decode('latin-1')
                encoding = 'latin-1'
                warnings.append("Used Latin-1 fallback encoding")
        
        # Language detection for code files
        language = self._detect_language(file_format, text)
        
        return ParsedFile(
            text=text,
            filename="",
            format=file_format,
            encoding=encoding,
            language=language,
            extraction_method="encoding detection + text read",
            warnings=warnings,
            metadata={
                "encoding_confidence": confidence,
                "detected_encoding": encoding_info.get('encoding'),
                "is_code_file": file_format not in ['txt', 'markdown']
            }
        )
    
    def _parse_markup(self, file_path: Path, file_format: str) -> ParsedFile:
        """Parse HTML/XML files using BeautifulSoup"""
        try:
            from bs4 import BeautifulSoup
            
            # Read file with encoding detection
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                
            encoding_info = chardet.detect(raw_data)
            encoding = encoding_info.get('encoding', 'utf-8')
            
            try:
                content = raw_data.decode(encoding)
            except:
                content = raw_data.decode('utf-8', errors='replace')
                encoding = 'utf-8'
            
            # Parse with BeautifulSoup
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "meta", "link"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            clean_text = '\n'.join(chunk for chunk in chunks if chunk)
            
            # Extract metadata
            metadata = {}
            if soup.title:
                metadata['title'] = soup.title.string
            
            meta_tags = soup.find_all('meta')
            for tag in meta_tags:
                if tag.get('name') and tag.get('content'):
                    metadata[tag.get('name')] = tag.get('content')
            
            return ParsedFile(
                text=clean_text,
                filename="",
                format=file_format,
                encoding=encoding,
                extraction_method="BeautifulSoup",
                metadata=metadata
            )
            
        except ImportError:
            raise ParsingError("HTML/XML parsing requires BeautifulSoup: pip install beautifulsoup4")
        except Exception as e:
            raise ParsingError(f"Markup parsing failed: {e}")
    
    def _parse_structured_data(self, file_path: Path, file_format: str) -> ParsedFile:
        """Parse structured data files (JSON, YAML, CSV, etc.)"""
        
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Format-specific parsing
        if file_format == 'json':
            text = self._parse_json_content(content)
        elif file_format == 'yaml':
            text = self._parse_yaml_content(content)
        elif file_format == 'csv':
            text = self._parse_csv_content(content)
        elif file_format == 'toml':
            text = self._parse_toml_content(content)
        elif file_format == 'ini':
            text = self._parse_ini_content(content)
        else:
            text = content  # Fallback to raw content
        
        return ParsedFile(
            text=text,
            filename="",
            format=file_format,
            encoding='utf-8',
            extraction_method=f"{file_format} parser",
            metadata={"original_format": file_format}
        )
    
    def _parse_json_content(self, content: str) -> str:
        """Parse JSON and convert to readable text"""
        try:
            import json
            data = json.loads(content)
            return self._format_structured_data(data, "JSON")
        except:
            return content  # Return raw if parsing fails
    
    def _parse_yaml_content(self, content: str) -> str:
        """Parse YAML and convert to readable text"""
        try:
            import yaml
            data = yaml.safe_load(content)
            return self._format_structured_data(data, "YAML")
        except:
            return content
    
    def _parse_csv_content(self, content: str) -> str:
        """Parse CSV and convert to readable text"""
        try:
            import csv
            import io
            
            reader = csv.reader(io.StringIO(content))
            rows = list(reader)
            
            if not rows:
                return content
                
            # Format as table
            formatted_rows = []
            for row in rows:
                formatted_rows.append(" | ".join(str(cell) for cell in row))
            
            return "\n".join(formatted_rows)
        except:
            return content
    
    def _parse_toml_content(self, content: str) -> str:
        """Parse TOML and convert to readable text"""
        try:
            import tomllib
            data = tomllib.loads(content)
            return self._format_structured_data(data, "TOML")
        except:
            return content
    
    def _parse_ini_content(self, content: str) -> str:
        """Parse INI and convert to readable text"""
        try:
            import configparser
            import io
            
            config = configparser.ConfigParser()
            config.read_string(content)
            
            sections = []
            for section_name in config.sections():
                section_items = []
                for key, value in config.items(section_name):
                    section_items.append(f"{key} = {value}")
                sections.append(f"[{section_name}]\n" + "\n".join(section_items))
            
            return "\n\n".join(sections)
        except:
            return content
    
    def _format_structured_data(self, data: Any, format_name: str) -> str:
        """Format structured data as readable text"""
        def format_item(item, indent=0):
            spaces = "  " * indent
            
            if isinstance(item, dict):
                lines = []
                for key, value in item.items():
                    if isinstance(value, (dict, list)):
                        lines.append(f"{spaces}{key}:")
                        lines.append(format_item(value, indent + 1))
                    else:
                        lines.append(f"{spaces}{key}: {value}")
                return "\n".join(lines)
            elif isinstance(item, list):
                lines = []
                for i, value in enumerate(item):
                    if isinstance(value, (dict, list)):
                        lines.append(f"{spaces}[{i}]:")
                        lines.append(format_item(value, indent + 1))
                    else:
                        lines.append(f"{spaces}- {value}")
                return "\n".join(lines)
            else:
                return f"{spaces}{item}"
        
        return format_item(data)
    
    def _detect_language(self, file_format: str, text: str) -> Optional[str]:
        """Detect programming language from format and content"""
        
        # Direct mapping from format
        language_map = {
            'python': 'Python',
            'javascript': 'JavaScript', 
            'typescript': 'TypeScript',
            'java': 'Java',
            'cpp': 'C++',
            'c': 'C',
            'go': 'Go',
            'rust': 'Rust',
            'ruby': 'Ruby',
            'php': 'PHP',
            'swift': 'Swift',
            'kotlin': 'Kotlin',
            'csharp': 'C#',
            'r': 'R',
            'scala': 'Scala',
            'shell': 'Shell',
            'markdown': 'Markdown'
        }
        
        return language_map.get(file_format)
    
    def _cleanup_temp_files(self):
        """Clean up any temporary files created during parsing"""
        for temp_file in self.temp_files:
            try:
                if temp_file.exists():
                    temp_file.unlink()
            except Exception as e:
                logger.warning(f"Failed to cleanup temp file {temp_file}: {e}")
        self.temp_files.clear()
    
    @classmethod
    def get_supported_formats(cls) -> List[str]:
        """Get list of all supported file formats"""
        return list(set(cls.SUPPORTED_FORMATS.values()))
    
    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """Check if a file format is supported"""
        suffix = Path(filename).suffix.lower()
        return suffix in cls.SUPPORTED_FORMATS


# Convenience functions for direct use
def parse_file(file_path: Union[str, Path], **kwargs) -> ParsedFile:
    """Quick file parsing function"""
    parser = FileParser(**kwargs)
    return parser.parse(file_path)


def get_supported_formats() -> List[str]:
    """Get list of supported formats"""
    return FileParser.get_supported_formats()


def is_supported_format(filename: str) -> bool:
    """Check if file format is supported"""
    return FileParser.is_supported(filename)